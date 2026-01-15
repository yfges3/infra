import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def is_korean(char):
    code = ord(char)
    return (
        0xAC00 <= code <= 0xD7AF or
        0x1100 <= code <= 0x11FF or
        0x3130 <= code <= 0x318F
    )

def calc_font_size(text, box_width_pt, box_height_pt):
    lines = text.split('\n')
    max_line_len = 0

    for line in lines:
        effective_len = 0
        for char in line:
            effective_len += 1.2 if is_korean(char) else 0.6
        max_line_len = max(max_line_len, effective_len)

    line_count = len(lines)
    line_spacing = 1.5

    min_font = 60  # 최소 폰트 크기 60pt
    for font_size in range(100, min_font - 1, -1):
        width_needed = max_line_len * font_size
        height_needed = line_count * font_size * line_spacing

        if width_needed <= box_width_pt and height_needed <= box_height_pt:
            return font_size
    return min_font

def create_ppt():
    prs = Presentation()
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 텍스트박스 크기 설정 (좌우 10px 패딩)
    padding_px = 10
    left = int(padding_px * 914400 / 96)  # px → EMU
    width = slide_width - left * 2

    box_width_pt = width / 12700  # EMU → 포인트
    box_height_pt = slide_height / 12700

    # txt 파일 목록 가져오기(정렬)
    file_names = sorted([file for file in os.listdir('.') if file.endswith('.txt')])

    # 1. 목차(첫) 슬라이드 생성
    toc_slide = prs.slides.add_slide(prs.slide_layouts[6])
    toc_slide.background.fill.solid()
    toc_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
    toc_box = toc_slide.shapes.add_textbox(left, 0, width, slide_height)
    toc_tf = toc_box.text_frame
    toc_tf.word_wrap = True
    toc_tf.vertical_anchor = MSO_ANCHOR.TOP

    # 2. 각 파일별 제목 슬라이드, 본문 슬라이드 생성 및 제목슬라이드 인덱스 기록
    title_slides = []
    title_slide_indices = []  # 목차에서 쓸 페이지 번호용
    for file in file_names:
        # 제목 슬라이드 생성
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        title_slides.append(title_slide)
        title_slide_indices.append(len(prs.slides))  # 1-based index (pptx는 0-based지만, 보통 1-based로 표시)

        title_slide.background.fill.solid()
        title_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
        title_box = title_slide.shapes.add_textbox(left, 0, width, slide_height)
        title_tf = title_box.text_frame
        title_tf.word_wrap = True
        title_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        title_p = title_tf.add_paragraph()
        title_p.text = file
        title_p.alignment = PP_ALIGN.CENTER
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 0)
        title_p.font.size = Pt(66)
        title_p.font.name = "a시네마m"

        # 본문 슬라이드들
        with open(file, 'r', encoding='utf-8') as f:
            content = f.read()
            paragraphs = [para.strip() for para in content.split('\n\n') if para.strip()]
            for para in paragraphs:
                content_slide = prs.slides.add_slide(prs.slide_layouts[6])
                content_slide.background.fill.solid()
                content_slide.background.fill.fore_color.rgb = RGBColor(0, 0, 0)
                textbox = content_slide.shapes.add_textbox(left, 0, width, slide_height)
                tf = textbox.text_frame
                tf.word_wrap = False
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.add_paragraph()
                p.alignment = PP_ALIGN.CENTER
                font_size = calc_font_size(para, box_width_pt, box_height_pt)
                p.text = para
                p.font.bold = True
                p.font.color.rgb = RGBColor(255, 255, 0)
                p.font.size = Pt(font_size)
                p.font.name = "a시네마m"

    # 3. 목차 슬라이드에 파일명, 페이지번호, 하이퍼링크 추가
    for idx, file in enumerate(file_names):
        p = toc_tf.add_paragraph()
        # 파일명과 페이지번호(제목슬라이드가 전체에서 몇 번째 슬라이드인지)
        page_num = title_slide_indices[idx]  # 1-based
        p.text = f"{file}  (p.{page_num})"
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = RGBColor(255, 255, 0)
        p.font.name = "a시네마m"
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0]
        run.hyperlink.address = None
        run.hyperlink.target_slide = title_slides[idx]

    prs.save('output.pptx')

if __name__ == "__main__":
    create_ppt()
